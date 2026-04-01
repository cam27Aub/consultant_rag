import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional


@dataclass
class PageUnit:
    source:     str
    doc_type:   str
    page:       int
    section:    str
    text:       str
    table_data: list = field(default_factory=list)   # raw table rows if present

    def to_dict(self) -> dict:
        return {
            "source":     self.source,
            "doc_type":   self.doc_type,
            "page":       self.page,
            "section":    self.section,
            "text":       self.text,
            "table_data": self.table_data,
        }


class DocumentCracker:
    """
    Cracks documents into PageUnit objects.
    Supports: .pdf, .pptx, .docx, .xlsx
    """

    def crack(self, filepath: Path) -> list[PageUnit]:
        ext = filepath.suffix.lower()
        dispatch = {
            ".pdf":  self._crack_pdf,
            ".pptx": self._crack_pptx,
            ".docx": self._crack_docx,
            ".xlsx": self._crack_xlsx,
            ".xls":  self._crack_xlsx,
        }
        fn = dispatch.get(ext)
        if not fn:
            print(f"  Unsupported file type: {ext} — skipping {filepath.name}")
            return []
        units = fn(filepath)
        print(f"  {filepath.name} → {len(units)} page/section units extracted")
        return units

    def _crack_pdf(self, filepath: Path) -> list[PageUnit]:
        import pdfplumber
        units = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                text = text.strip()

                table_rows = []
                for table in page.extract_tables():
                    if not table:
                        continue
                    for row in table:
                        row_text = " | ".join(
                            str(cell).strip() for cell in row if cell and str(cell).strip()
                        )
                        if row_text:
                            table_rows.append(row_text)

                combined = text
                if table_rows:
                    combined += "\n\nTABLE DATA:\n" + "\n".join(table_rows)

                if len(combined.split()) < 10:
                    continue

                first_line = next(
                    (ln.strip() for ln in combined.split("\n") if ln.strip()), f"Page {i+1}"
                )
                section = first_line[:80] if len(first_line) <= 80 else f"Page {i+1}"

                units.append(PageUnit(
                    source=filepath.name,
                    doc_type="pdf",
                    page=i + 1,
                    section=section,
                    text=combined,
                    table_data=table_rows,
                ))
        return units

    def _crack_pptx(self, filepath: Path) -> list[PageUnit]:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation(filepath)
        units = []

        for i, slide in enumerate(prs.slides):
            texts = []
            title_text = ""
            table_rows = []

            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        row_str = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_str:
                            table_rows.append(row_str)
                    texts.extend(table_rows)
                    continue

                if not shape.has_text_frame:
                    continue

                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if not line:
                        continue
                    try:
                        if (shape.placeholder_format and
                                shape.placeholder_format.idx in (0, 1)):
                            title_text = line
                    except Exception:
                        pass
                    texts.append(line)

            combined = "\n".join(texts).strip()
            if len(combined.split()) < 8:
                continue

            units.append(PageUnit(
                source=filepath.name,
                doc_type="pptx",
                page=i + 1,
                section=title_text or f"Slide {i+1}",
                text=combined,
                table_data=table_rows,
            ))
        return units

    def _crack_docx(self, filepath: Path) -> list[PageUnit]:
        from docx import Document

        doc = Document(filepath)
        units = []
        current_section = "Introduction"
        current_lines = []
        current_tables = []
        section_index = 0

        def flush():
            nonlocal current_lines, current_tables, section_index
            if not current_lines and not current_tables:
                return
            body = "\n".join(current_lines)
            if current_tables:
                body += "\n\nTABLE DATA:\n" + "\n".join(current_tables)
            if len(body.split()) >= 8:
                section_index += 1
                units.append(PageUnit(
                    source=filepath.name,
                    doc_type="docx",
                    page=section_index,
                    section=current_section,
                    text=body.strip(),
                    table_data=list(current_tables),
                ))
            current_lines.clear()
            current_tables.clear()

        for element in doc.element.body:
            tag = element.tag.split("}")[-1]

            if tag == "p":
                from docx.oxml.ns import qn
                style_name = ""
                pPr = element.find(qn("w:pPr"))
                if pPr is not None:
                    pStyle = pPr.find(qn("w:pStyle"))
                    if pStyle is not None:
                        style_name = pStyle.get(qn("w:val"), "")

                text = "".join(
                    node.text for node in element.iter()
                    if node.tag.endswith("}t") and node.text
                ).strip()

                if not text:
                    continue

                if "Heading" in style_name or style_name.startswith("heading"):
                    flush()
                    current_section = text
                else:
                    current_lines.append(text)

            elif tag == "tbl":
                from docx.oxml.ns import qn
                rows = element.findall(".//" + qn("w:tr"))
                for row in rows:
                    cells = row.findall(".//" + qn("w:tc"))
                    cell_texts = []
                    for cell in cells:
                        ct = "".join(
                            n.text for n in cell.iter()
                            if n.tag.endswith("}t") and n.text
                        ).strip()
                        if ct:
                            cell_texts.append(ct)
                    if cell_texts:
                        current_tables.append(" | ".join(cell_texts))

        flush()
        return units

    def _crack_xlsx(self, filepath: Path) -> list[PageUnit]:
        import openpyxl

        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        units = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            header_row = None

            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                vals = [str(v).strip() for v in row if v is not None and str(v).strip() and str(v).strip() != "None"]
                if not vals:
                    continue
                if row_idx == 0:
                    header_row = vals
                rows_text.append(" | ".join(vals))

            if not rows_text:
                continue

            # Build a prose-friendly summary line
            summary = f"Sheet: {sheet_name}."
            if header_row:
                summary += f" Columns: {', '.join(header_row[:8])}."

            full_text = summary + "\n\n" + "\n".join(rows_text)

            units.append(PageUnit(
                source=filepath.name,
                doc_type="xlsx",
                page=wb.sheetnames.index(sheet_name) + 1,
                section=sheet_name,
                text=full_text,
                table_data=rows_text,
            ))

        wb.close()
        return units
